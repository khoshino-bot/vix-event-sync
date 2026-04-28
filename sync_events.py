#!/usr/bin/env python3
"""
イベント起案・振り返り自動チェックスクリプト（無料版）

検証ロジック:
  ① 件名から店舗・日付を正規表現で抽出
  ② 添付ファイル(pptx/xlsx)から店舗・日付を正規表現で抽出
  ① == ② → シート書き込み
  ① != ② → 送信者に通知メール（シートは更新しない）
  添付なし → 件名のみで処理
"""

import os, re, json, base64, tempfile, time, unicodedata
from datetime import datetime, date, timedelta
from collections import Counter
from email.mime.text import MIMEText
import email.utils as email_utils
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from pptx import Presentation
import openpyxl

# ===== 設定 =====
SPREADSHEET_ID = os.environ.get("SPREADSHEET_ID", "1ia4DThYgqZ3WdyeoQcBpKj5c_t8bJdYgFZIY7Ge2ubE")
PROCESSED_FILE = "processed_ids.json"
SENDER_DOMAIN  = os.environ.get("SENDER_DOMAIN", "vix.co.jp")
# 相違通知メールの送信（デフォルト有効）
NOTIFY_EMAIL   = os.environ.get("NOTIFY_EMAIL", "true").lower() == "true"
# 相違通知メールの CC 先（管理者）
OWNER_EMAIL    = os.environ.get("OWNER_EMAIL", "k.hoshino@vix.co.jp")
# RESET_PROCESSED=true で処理済みキャッシュを無視（シート再投入時に使用）
RESET_PROCESSED = os.environ.get("RESET_PROCESSED", "false").lower() == "true"
# FILTER_YYYYMM=202604 で対象年月を指定。未設定時は当月を自動検出。
_fym = os.environ.get("FILTER_YYYYMM", "").strip()
if len(_fym) == 6:
    FILTER_YM = (int(_fym[:4]), int(_fym[4:]))
else:
    _now = datetime.now()
    FILTER_YM = (_now.year, _now.month)

STORE_MAP = {
    "経堂":     ["経堂"],
    "日暮里":   ["日暮里"],
    "町田":     ["町田", "コビルナ"],
    "ゆめが丘": ["ゆめが丘", "ゆめがおか", "ゆめが丘ソラトス"],
    "立川":     ["立川"],
    "三軒茶屋": ["三軒茶屋"],
    "竹ノ塚":   ["竹ノ塚"],
    "木の葉":   ["木の葉", "橋本"],
}

# 店舗の表示順（シートでの並び順）
STORE_ORDER = ["竹ノ塚", "三軒茶屋", "日暮里", "町田", "木の葉", "立川", "経堂", "ゆめが丘"]

# ===== KPI スプレッドシート設定 =====
KPI_SPREADSHEET_ID = os.environ.get("KPI_SPREADSHEET_ID", "")
KPI_SHEET_NAME     = os.environ.get("KPI_SHEET_NAME", "貼り付け")

# シートのD列（店舗名）→ KPI C列（楽天店舗名）のマッピング
STORE_TO_KPI = {
    "三軒茶屋": "楽天三茶",
    "日暮里":   "楽天日暮里",
    "町田":     "楽天町田",
    "木の葉":   "楽天木の葉ﾓｰﾙ",
    "立川":     "楽天立川",
    "経堂":     "楽天経堂",
    "ゆめが丘": "楽天ゆめが丘",
    "竹ノ塚":   "楽天竹ノ塚",
}

def store_sort_key(store):
    """店舗名を並び順のインデックスに変換（未知の店舗は末尾）"""
    try:
        return STORE_ORDER.index(store)
    except ValueError:
        return len(STORE_ORDER)

# ===== 認証 =====
def build_creds(prefix):
    creds = Credentials(
        token=None,
        refresh_token=os.environ[f"{prefix}_REFRESH_TOKEN"],
        token_uri="https://oauth2.googleapis.com/token",
        client_id=os.environ[f"{prefix}_CLIENT_ID"],
        client_secret=os.environ[f"{prefix}_CLIENT_SECRET"],
    )
    creds.refresh(Request())
    return creds

# ===== 処理済みID =====
def load_processed():
    return set(json.load(open(PROCESSED_FILE))) if os.path.exists(PROCESSED_FILE) else set()

def save_processed(ids):
    json.dump(list(ids), open(PROCESSED_FILE, "w"), indent=2)

# ===== 情報抽出（正規表現）=====
def nkfc(text: str) -> str:
    """全角英数字・記号を半角に正規化（全角スラッシュ・チルダ等も対応）"""
    return unicodedata.normalize('NFKC', text)

def normalize_store(text):
    text = nkfc(text)
    for store, keywords in STORE_MAP.items():
        if any(kw in text for kw in keywords):
            return store
    return None

def _expand_range(start: date, end: date, max_days=31) -> list:
    """start〜end を日単位に展開（最大max_days日）"""
    if start > end:
        return []
    result = []
    d = start
    while d <= end and len(result) < max_days:
        result.append(d)
        d += timedelta(days=1)
    return result

def extract_dates(text):
    """
    テキストから日付を抽出する（範囲・コンテキスト年対応版）。

    優先順位:
      ① YYYYMMDD / YYYYMMDD.DD範囲（明示年）
      ② YYYY年M月D日〜M月D日 範囲（明示年）
      ③ YYYY年M月D日 単発（明示年）
      以上を「explicit」として確定し、その(月,日)ペアをロック。
      コンテキスト年 = explicit 内の最頻年（なければ今年）

      ④ M/D〜M/D 範囲（コンテキスト年）
      ⑤ M月D日〜M月D日 範囲（コンテキスト年）
      ⑥ M月D〜D日 同月範囲（コンテキスト年）
      ⑦ M月D,D,D日 / M月D.D.D日 リスト（コンテキスト年）
      ⑧ M/D 単発（コンテキスト年）
      ④〜⑧を「implicit」として、(月,日)が explicit と重複しなければ追加。
    """
    cur_year = datetime.now().year
    # 全角数字・スラッシュ・チルダ等を半角に統一してからパース
    text = nkfc(text)
    explicit = set()
    SEP = r'[〜~\-]'  # 範囲区切り文字（NFKC後: 〜 のみ全角、他は半角）

    # ① YYYY-MM-DD 形式（openpyxl が date/datetime セルを str() 変換した結果）
    # 例: "2026-04-11" or "2026-04-11 00:00:00"
    for m in re.finditer(r'(20\d{2})-(\d{2})-(\d{2})', text):
        try:
            explicit.add(date(int(m.group(1)), int(m.group(2)), int(m.group(3))))
        except ValueError:
            pass

    # ① YYYYMMDD 単発
    for m in re.finditer(r'(20\d{2})(\d{2})(\d{2})', text):
        try:
            explicit.add(date(int(m.group(1)), int(m.group(2)), int(m.group(3))))
        except ValueError:
            pass

    # ① YYYYMMDD.DD~DD.DD → 単発 + 範囲 + 単発 (例: 20260408.11~14.18 → 4/8単発, 4/11-14範囲, 4/18単発)
    for m in re.finditer(r'(20\d{2})(\d{2})(\d{2})[.](\d{1,2})' + SEP + r'(\d{1,2})[.](\d{1,2})', text):
        yr, mo = int(m.group(1)), int(m.group(2))
        try:
            explicit.add(date(yr, mo, int(m.group(3))))   # DD1: 単発
            for d in _expand_range(date(yr, mo, int(m.group(4))), date(yr, mo, int(m.group(5)))):
                explicit.add(d)                            # DD2〜DD3: 範囲
            explicit.add(date(yr, mo, int(m.group(6))))   # DD4: 単発
        except ValueError:
            pass

    # ① YYYYMMDD.DD → 同月の2つの個別日付 (例: 20260404.25 → 4/4 と 4/25)
    # ※ 範囲ではなく「日付1.日付2」という列挙形式
    for m in re.finditer(r'(20\d{2})(\d{2})(\d{2})[.](\d{1,2})', text):
        yr, mo, d1, d2 = int(m.group(1)), int(m.group(2)), int(m.group(3)), int(m.group(4))
        try:
            explicit.add(date(yr, mo, d1))
            explicit.add(date(yr, mo, d2))
        except ValueError:
            pass

    # ② YYYY年M月D日〜M月D日 範囲（日あり形式）
    for m in re.finditer(r'(20\d{2})年(\d{1,2})月(\d{1,2})日' + SEP + r'(\d{1,2})月(\d{1,2})日', text):
        yr = int(m.group(1))
        try:
            for d in _expand_range(date(yr, int(m.group(2)), int(m.group(3))),
                                   date(yr, int(m.group(4)), int(m.group(5)))):
                explicit.add(d)
        except ValueError:
            pass

    # ② YYYY年M月D〜M月D日 範囲（始点の「日」なし形式 例: 2026年4月16〜4月27日）
    for m in re.finditer(r'(20\d{2})年(\d{1,2})月(\d{1,2})' + SEP + r'(\d{1,2})月(\d{1,2})日', text):
        yr = int(m.group(1))
        try:
            for d in _expand_range(date(yr, int(m.group(2)), int(m.group(3))),
                                   date(yr, int(m.group(4)), int(m.group(5)))):
                explicit.add(d)
        except ValueError:
            pass

    # ② YYYY年M月D日〜D日 同月範囲（終点の月なし 例: 2025年12月9日〜12日）
    for m in re.finditer(r'(20\d{2})年(\d{1,2})月(\d{1,2})日' + SEP + r'(\d{1,2})日', text):
        yr, mo = int(m.group(1)), int(m.group(2))
        try:
            for d in _expand_range(date(yr, mo, int(m.group(3))),
                                   date(yr, mo, int(m.group(4)))):
                explicit.add(d)
        except ValueError:
            pass

    # ② YYYY年M月D〜D日 同月範囲（始点の「日」なし 例: 2026年4月16〜27日）
    for m in re.finditer(r'(20\d{2})年(\d{1,2})月(\d{1,2})' + SEP + r'(\d{1,2})日', text):
        yr, mo = int(m.group(1)), int(m.group(2))
        try:
            for d in _expand_range(date(yr, mo, int(m.group(3))),
                                   date(yr, mo, int(m.group(4)))):
                explicit.add(d)
        except ValueError:
            pass

    # ② Extended: YYYY年M月 アンカーから前方スキャンして D日〜D日 / D日 を収集
    # 例: "2025年12月9日〜12日、22日〜26日" → 9〜12 と 22〜26 を両方取得
    for m in re.finditer(r'(20\d{2})年(\d{1,2})月', text):
        yr, mo = int(m.group(1)), int(m.group(2))
        chunk = text[m.end():m.end() + 120]
        # 次の別月（M月）が現れたらそこで打ち切り
        next_mo = re.search(r'\d+月', chunk)
        if next_mo:
            chunk = chunk[:next_mo.start()]
        # D日〜D日 範囲（両端に日あり）
        for rm in re.finditer(r'(\d{1,2})日' + SEP + r'(\d{1,2})日', chunk):
            try:
                for d in _expand_range(date(yr, mo, int(rm.group(1))),
                                       date(yr, mo, int(rm.group(2)))):
                    explicit.add(d)
            except ValueError:
                pass
        # D〜D日 範囲（終端のみ日あり 例: 6~10日, 2~3日）
        for rm in re.finditer(r'(\d{1,2})' + SEP + r'(\d{1,2})日', chunk):
            try:
                for d in _expand_range(date(yr, mo, int(rm.group(1))),
                                       date(yr, mo, int(rm.group(2)))):
                    explicit.add(d)
            except ValueError:
                pass
        # D日 単発
        for dm in re.finditer(r'(\d{1,2})日', chunk):
            try:
                explicit.add(date(yr, mo, int(dm.group(1))))
            except ValueError:
                pass

    # ③ YYYY年M月D日 単発
    for m in re.finditer(r'(20\d{2})年(\d{1,2})月(\d{1,2})日', text):
        try:
            explicit.add(date(int(m.group(1)), int(m.group(2)), int(m.group(3))))
        except ValueError:
            pass

    # コンテキスト年: テキスト内の明示年の最頻値（なければ今年）
    explicit_years = [int(y) for y in re.findall(r'(20\d{2})年', text)]
    ctx_year = Counter(explicit_years).most_common(1)[0][0] if explicit_years else cur_year

    explicit_md = {(d.month, d.day) for d in explicit}
    implicit = set()

    def add_implicit(mo, dy):
        if 1 <= mo <= 12 and 1 <= dy <= 31 and (mo, dy) not in explicit_md:
            try:
                implicit.add(date(ctx_year, mo, dy))
            except ValueError:
                pass

    def add_range_implicit(mo1, d1, mo2, d2):
        try:
            for d in _expand_range(date(ctx_year, mo1, d1), date(ctx_year, mo2, d2)):
                if (d.month, d.day) not in explicit_md:
                    implicit.add(d)
        except ValueError:
            pass

    # ④ M/D〜M/D 範囲 (例: 4/12~4/16)
    for m in re.finditer(r'(\d{1,2})/(\d{1,2})' + SEP + r'(\d{1,2})/(\d{1,2})', text):
        add_range_implicit(int(m.group(1)), int(m.group(2)),
                           int(m.group(3)), int(m.group(4)))

    # ⑤ M月D日〜M月D日 範囲（年なし）
    for m in re.finditer(r'(\d{1,2})月(\d{1,2})日' + SEP + r'(\d{1,2})月(\d{1,2})日', text):
        add_range_implicit(int(m.group(1)), int(m.group(2)),
                           int(m.group(3)), int(m.group(4)))

    # ⑥ M月D〜D日 同月範囲 (例: 12月8〜12日, 4/16~27日)
    for m in re.finditer(r'(\d{1,2})月(\d{1,2})' + SEP + r'(\d{1,2})日', text):
        mo = int(m.group(1))
        d1, d2 = int(m.group(2)), int(m.group(3))
        add_range_implicit(mo, d1, mo, d2)

    # ⑦ M月D,D,D日 / M月D.D.D日 リスト (例: 4月3,6,9日 / 4月16.17日)
    for m in re.finditer(r'(\d{1,2})月([\d,、・.\s]+)日', text):
        mo = int(m.group(1))
        if not (1 <= mo <= 12):
            continue
        for day_str in re.split(r'[,、・.\s]+', m.group(2)):
            day_str = day_str.strip()
            if day_str:
                try:
                    add_implicit(mo, int(day_str))
                except ValueError:
                    pass

    # ⑧ M/D 単発 (例: 4/5)
    for m in re.finditer(r'(\d{1,2})/(\d{1,2})', text):
        add_implicit(int(m.group(1)), int(m.group(2)))

    return sorted(explicit | implicit)

# ===== 添付ファイル読み取り =====
def extract_text_from_pptx(path):
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
    return "\n".join(parts)

def extract_text_from_xlsx(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)

def get_attachment_text(gmail_svc, msg_id, payload):
    texts = []
    def walk(parts):
        for part in parts:
            fname = part.get("filename", "")
            att_id = part["body"].get("attachmentId", "")
            if not att_id:
                if "parts" in part:
                    walk(part["parts"])
                continue
            att = gmail_svc.users().messages().attachments().get(
                userId="me", messageId=msg_id, id=att_id
            ).execute()
            data = base64.urlsafe_b64decode(att["data"])
            ext = os.path.splitext(fname)[1].lower()
            tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
            tmp.write(data)
            tmp.close()
            try:
                if fname.endswith(".pptx"):
                    texts.append(extract_text_from_pptx(tmp.name))
                elif fname.endswith((".xlsx", ".xls")):
                    texts.append(extract_text_from_xlsx(tmp.name))
            finally:
                os.unlink(tmp.name)
            if "parts" in part:
                walk(part["parts"])
    walk(payload.get("parts", []))
    return "\n".join(texts)

# ===== 二重検証（件名 vs 添付）=====
def verify(subject, att_text):
    """
    件名と添付ファイルのテキストを正規表現で比較。
    Returns: (ok, store, dates, diff_msg)
    """
    s_store = normalize_store(subject)
    s_dates = extract_dates(subject)
    a_store = normalize_store(att_text)
    a_dates = extract_dates(att_text)

    print(f"    [件名]   store={s_store} dates={s_dates}")
    print(f"    [添付]   store={a_store} dates={a_dates}")

    # 添付から日付が取れない場合は件名の日付を信頼する
    # （添付が未対応フォーマット or 日付セルのみ → 誤検知を防ぐ）
    if not a_dates:
        # 店舗名が両方あって食い違う場合だけ不一致とする
        if a_store and s_store and a_store != s_store:
            return False, a_store, s_dates, f"店舗: 件名={s_store} / 添付={a_store}"
        return True, a_store or s_store, s_dates, ""

    diffs = []
    if s_store and a_store and s_store != a_store:
        diffs.append(f"店舗: 件名={s_store} / 添付={a_store}")

    s_set, a_set = set(s_dates), set(a_dates)
    only_subject = sorted(s_set - a_set)
    only_attach  = sorted(a_set - s_set)
    if only_subject:
        diffs.append(f"件名のみの日付: {only_subject}")
    if only_attach:
        diffs.append(f"添付のみの日付: {only_attach}")

    if diffs:
        return False, a_store or s_store, sorted(a_set | s_set), "\n".join(diffs)

    # 一致 → 添付の情報を正とする
    return True, a_store or s_store, sorted(a_set | s_set), ""

# ===== 通知メール =====
def send_discrepancy_email(gmail_svc, original_msg, subject, diff_msg):
    headers = {h["name"]: h["value"] for h in original_msg["payload"]["headers"]}
    raw_from = headers.get("From", "")
    # 表示名付き "名前 <email@domain>" からメールアドレスのみ抽出
    _, addr = email_utils.parseaddr(raw_from)
    to   = addr if addr else raw_from
    subj = "Re: " + headers.get("Subject", "")
    body = f"""お疲れ様です。
イベント申請の件名と添付ファイルの内容が合っていないようです。確認していただけますか？

【気になった箇所】
{diff_msg}

問題なければそのままで大丈夫です。修正が必要な場合は再送をお願いします。
再送する場合は件名の前に「再送」をつけてください。

---
このメールは自動送信です。
"""
    msg_obj = MIMEText(body, "plain", "utf-8")
    msg_obj["To"]      = to
    msg_obj["Subject"] = subj
    # 管理者（OWNER_EMAIL）を CC に追加
    if OWNER_EMAIL and OWNER_EMAIL != to:
        msg_obj["Cc"] = OWNER_EMAIL
    encoded = base64.urlsafe_b64encode(msg_obj.as_bytes()).decode()
    gmail_svc.users().messages().send(
        userId="me", body={"raw": encoded}
    ).execute()
    cc_info = f" / CC: {OWNER_EMAIL}" if OWNER_EMAIL and OWNER_EMAIL != to else ""
    print(f"    → 通知メール送信: {to}{cc_info}")

# ===== シート操作 =====
def get_current_sheet_name():
    # TARGET_SHEET 環境変数が設定されていればそれを優先
    override = os.environ.get("TARGET_SHEET", "").strip()
    if override:
        return override
    now = datetime.now()
    return f"{str(now.year)[2:]}{str(now.month).zfill(2)}"

def load_sheet_rows(sheets_svc, sheet_name):
    result = sheets_svc.spreadsheets().values().get(
        spreadsheetId=SPREADSHEET_ID,
        range=f"'{sheet_name}'!A1:Z300"
    ).execute()
    rows = result.get("values", [])
    row_map   = {}
    date_rows = []
    for i, row in enumerate(rows):
        if i < 3:
            continue
        store    = row[3].strip() if len(row) > 3 else ""
        date_str = row[4].strip() if len(row) > 4 else ""
        if not store or not date_str:
            continue
        try:
            d = datetime.strptime(date_str, "%Y/%m/%d").date()
            row_map[(store, d)] = i + 1
            date_rows.append((store_sort_key(store), d, i + 1))
        except ValueError:
            pass
    return row_map, sorted(date_rows)

def deduplicate_sheet(sheets_svc, sheet_name, sheet_id):
    """シート内の重複（店舗・日付）行を検出して削除する（先に出た行を削除＝再送の上書きを優先）"""
    result = execute_with_retry(sheets_svc.spreadsheets().values().get(
        spreadsheetId=SPREADSHEET_ID,
        range=f"'{sheet_name}'!A1:Z300"
    ))
    rows = result.get("values", [])
    seen = {}
    to_delete = []
    for i, row in enumerate(rows):
        if i < 3:
            continue
        store    = row[3].strip() if len(row) > 3 else ""
        date_str = row[4].strip() if len(row) > 4 else ""
        if not store or not date_str:
            continue
        key = (store, date_str)
        if key in seen:
            to_delete.append(seen[key])  # 先に出た方（古い行）を削除
            seen[key] = i + 1            # 最新行番号に更新
        else:
            seen[key] = i + 1
    if not to_delete:
        return
    print(f"★ 重複行を削除: {len(to_delete)}行 {sorted(to_delete)}")
    # 後ろから削除して行番号ずれを防ぐ
    for row_num in sorted(to_delete, reverse=True):
        execute_with_retry(sheets_svc.spreadsheets().batchUpdate(
            spreadsheetId=SPREADSHEET_ID,
            body={"requests": [{"deleteDimension": {
                "range": {
                    "sheetId":    sheet_id,
                    "dimension":  "ROWS",
                    "startIndex": row_num - 1,
                    "endIndex":   row_num,
                }
            }}]}
        ))


def delete_rows_by_dates(sheets_svc, sheet_name, sheet_id, target_dates, row_map):
    """
    再送メール処理用: 指定日付に一致する既存行をすべて削除する。
    同じ日付であれば店舗名が変わっていても古い行を除去できる。
    Returns: 削除があった場合 True
    """
    target_set = set(target_dates)
    rows_to_delete = sorted(
        [row_num for (_, d), row_num in row_map.items() if d in target_set],
        reverse=True,
    )
    if not rows_to_delete:
        return False
    print(f"    [再送] 既存行を削除: {len(rows_to_delete)}行 {rows_to_delete}")
    for row_num in rows_to_delete:
        execute_with_retry(sheets_svc.spreadsheets().batchUpdate(
            spreadsheetId=SPREADSHEET_ID,
            body={"requests": [{"deleteDimension": {
                "range": {
                    "sheetId":    sheet_id,
                    "dimension":  "ROWS",
                    "startIndex": row_num - 1,
                    "endIndex":   row_num,
                }
            }}]}
        ))
    return True

def execute_with_retry(request, max_retries=6, pace_sec=1.1):
    """
    429/500/503 に対して指数バックオフでリトライ。
    成功時も pace_sec 待機してレートリミット(60req/min)を遵守。
    """
    for attempt in range(max_retries):
        try:
            result = request.execute()
            time.sleep(pace_sec)   # 成功後もペーシング（約54req/min < 60の上限）
            return result
        except HttpError as e:
            if e.resp.status in (429, 500, 503) and attempt < max_retries - 1:
                wait = 2 ** (attempt + 1)  # 2, 4, 8, 16, 32, 64 秒
                print(f"    [リトライ {attempt+1}/{max_retries-1}] HTTP {e.resp.status} → {wait}秒待機")
                time.sleep(wait)
            else:
                raise

def get_sheet_id(sheets_svc, sheet_name):
    meta = execute_with_retry(sheets_svc.spreadsheets().get(spreadsheetId=SPREADSHEET_ID))
    return next(s["properties"]["sheetId"] for s in meta["sheets"]
                if s["properties"]["title"] == sheet_name)

def ensure_current_sheet(sheets_svc, sheet_name: str) -> int:
    """
    当月シートが存在しない場合はテンプレートシートを複製して作成する。
    既に存在する場合はそのまま sheetId を返す（冪等）。
    """
    meta   = execute_with_retry(sheets_svc.spreadsheets().get(spreadsheetId=SPREADSHEET_ID))
    sheets = meta["sheets"]
    titles = {s["properties"]["title"]: s for s in sheets}

    if sheet_name in titles:
        return titles[sheet_name]["properties"]["sheetId"]

    # シートが存在しない → テンプレートから自動作成
    tpl_name = os.environ.get("TEMPLATE_SHEET_NAME", "テンプレート")
    if tpl_name not in titles:
        raise SystemExit(
            f"[ERROR] シート '{sheet_name}' が存在せず、"
            f"テンプレートシート '{tpl_name}' も見つかりません。\n"
            f"スプレッドシート内に '{tpl_name}' シートを作成してください。"
        )

    print(f"シート '{sheet_name}' が存在しないため '{tpl_name}' から作成します...")
    yymm_sheets = sorted(
        [s for s in sheets
         if len(s["properties"]["title"]) == 4 and s["properties"]["title"].isdigit()],
        key=lambda s: s["properties"]["title"]
    )
    insert_index = (yymm_sheets[-1]["properties"]["index"] + 1) if yymm_sheets else 1

    res = execute_with_retry(sheets_svc.spreadsheets().batchUpdate(
        spreadsheetId=SPREADSHEET_ID,
        body={"requests": [{"duplicateSheet": {
            "sourceSheetId":    titles[tpl_name]["properties"]["sheetId"],
            "insertSheetIndex": insert_index,
            "newSheetName":     sheet_name,
        }}]}
    ))
    new_id = res["replies"][0]["duplicateSheet"]["properties"]["sheetId"]
    print(f"シート '{sheet_name}' を作成しました (sheetId={new_id})")
    return new_id

def insert_event_row(sheets_svc, sheet_name, store, d, date_rows, row_map, sheet_id,
                     first_formatted_row=4):
    """
    重複チェック付き行挿入。
    - (store, d) が既に row_map に存在する場合はスキップ
    - inheritFromBefore の選択:
        insert_at > first_formatted_row → True（上の行から書式継承）
        insert_at <= first_formatted_row → False（下の行=例行から書式継承）
      これにより例行より前の日付でもドロップダウン書式が正しく引き継がれる
    """
    if (store, d) in row_map:
        print(f"    → スキップ(既存): {store} {d}")
        return date_rows, row_map

    # 店舗順・日付順の挿入位置を決定
    new_key = (store_sort_key(store), d)
    insert_at = 4
    for sk, dt, row_num in date_rows:
        if new_key < (sk, dt):
            insert_at = row_num
            break
        insert_at = row_num + 1

    # 例行より後 → 上から継承、例行以前 → 下（例行）から継承
    inherit_from_before = (insert_at > first_formatted_row)

    # 行挿入後のテンプレート行（0-indexed）:
    #   inheritFromBefore=True  → 上の行 (insert_at-2)  ← 挿入で位置変わらず
    #   inheritFromBefore=False → 下の例行 (insert_at)  ← 挿入で+1シフト済み
    template_row_0 = (insert_at - 2) if inherit_from_before else insert_at
    new_row_0      = insert_at - 1

    # M:T列(12-19) と V:W列(21-22) の数式を明示的にコピー（数式のセル参照を新行に合わせて調整）
    formula_ranges = [(12, 20), (21, 23)]  # (startColumnIndex, endColumnIndex) 0-indexed
    copy_requests = [
        {"copyPaste": {
            "source": {
                "sheetId": sheet_id,
                "startRowIndex": template_row_0, "endRowIndex": template_row_0 + 1,
                "startColumnIndex": c_start,     "endColumnIndex": c_end,
            },
            "destination": {
                "sheetId": sheet_id,
                "startRowIndex": new_row_0,      "endRowIndex": new_row_0 + 1,
                "startColumnIndex": c_start,     "endColumnIndex": c_end,
            },
            "pasteType": "PASTE_FORMULA",
            "pasteOrientation": "NORMAL",
        }}
        for c_start, c_end in formula_ranges
    ]

    execute_with_retry(sheets_svc.spreadsheets().batchUpdate(
        spreadsheetId=SPREADSHEET_ID,
        body={"requests": [
            {"insertDimension": {
                "range": {
                    "sheetId":    sheet_id,
                    "dimension":  "ROWS",
                    "startIndex": insert_at - 1,
                    "endIndex":   insert_at,
                },
                "inheritFromBefore": inherit_from_before
            }},
            *copy_requests,
        ]}
    ))

    execute_with_retry(sheets_svc.spreadsheets().values().update(
        spreadsheetId=SPREADSHEET_ID,
        range=f"'{sheet_name}'!B{insert_at}:E{insert_at}",
        valueInputOption="RAW",
        body={"values": [["〇", "", store, d.strftime("%Y/%m/%d")]]}
    ))
    print(f"    → 行{insert_at}に挿入: {store} {d} 起案〇")

    # メモリ上で date_rows と row_map を更新（行挿入で以降の行番号が+1ずれる）
    new_date_rows = [(sk, dt, rn + 1) if rn >= insert_at else (sk, dt, rn) for sk, dt, rn in date_rows]
    new_date_rows.append((store_sort_key(store), d, insert_at))
    new_row_map   = {(s, dt): (rn + 1 if rn >= insert_at else rn) for (s, dt), rn in row_map.items()}
    new_row_map[(store, d)] = insert_at
    return sorted(new_date_rows), new_row_map

def update_cell(sheets_svc, sheet_name, row_num, col):
    col_letter = "B" if col == "起案" else "C"
    cell = f"'{sheet_name}'!{col_letter}{row_num}"
    cur  = execute_with_retry(sheets_svc.spreadsheets().values().get(
        spreadsheetId=SPREADSHEET_ID, range=cell
    ))
    if cur.get("values", [[""]])[0][0] == "〇":
        return False
    execute_with_retry(sheets_svc.spreadsheets().values().update(
        spreadsheetId=SPREADSHEET_ID, range=cell,
        valueInputOption="RAW", body={"values": [["〇"]]}
    ))
    return True

# ===== KPI列（X〜AA）自動入力 =====
def parse_kpi_date(date_str: str):
    """KPI B列の日付文字列を date に変換（"26/04/03" or "26/4/3"）"""
    m = re.match(r'(\d{2})/(\d{1,2})/(\d{1,2})', (date_str or "").strip())
    if not m:
        return None
    try:
        return date(2000 + int(m.group(1)), int(m.group(2)), int(m.group(3)))
    except ValueError:
        return None

def categorize_kpi_row(j_val: str) -> dict:
    """J列の商品名をカテゴリに分類して {mnp, hikari, turbo, card} を返す"""
    j = (j_val or "").strip()
    return {
        "mnp":    1 if ("MNP" in j or "新規/ﾏｲｸﾞﾚ" in j) else 0,
        "hikari": 1 if "光" in j else 0,
        "turbo":  1 if "楽天Turbo" in j else 0,
        "card":   1 if (("楽天ｶｰﾄﾞ" in j or "楽天ｺﾞｰﾙﾄﾞｶｰﾄﾞ" in j)
                         and "審査待ち" not in j and "落ち" not in j) else 0,
    }

def fill_kpi_columns(sheets_svc, event_sheet_name: str):
    """
    X〜AA列: イベント日と一致したｸﾛｰｻﾞｰ実績を各イベント行に書き込む。
    AD〜AH列: イベント日以外のｸﾛｰｻﾞｰ獲得数（戻り）を店舗ごとに集計して固定位置に書き込む。
      行1=セクションタイトル、行2=ヘッダー、行3〜10=STORE_ORDER順の店舗別合計。
    """
    print("\n=== KPI列入力開始 ===")

    # ① KPI 貼り付けシート全行取得
    kpi_rows = execute_with_retry(sheets_svc.spreadsheets().values().get(
        spreadsheetId=KPI_SPREADSHEET_ID,
        range=f"'{KPI_SHEET_NAME}'!A1:J10000"
    )).get("values", [])
    print(f"KPI総行数: {len(kpi_rows)}")

    # ② ｸﾛｰｻﾞｰ行を (date, kpi_store) キーで集計
    kpi_counts: dict = {}
    for row in kpi_rows:
        d_role = row[3].strip() if len(row) > 3 else ""
        if d_role != "ｸﾛｰｻﾞｰ":
            continue
        d = parse_kpi_date(row[1] if len(row) > 1 else "")
        if not d:
            continue
        c_store = row[2].strip() if len(row) > 2 else ""
        key = (d, c_store)
        if key not in kpi_counts:
            kpi_counts[key] = {"mnp": 0, "hikari": 0, "turbo": 0, "card": 0}
        cats = categorize_kpi_row(row[9] if len(row) > 9 else "")
        for k in cats:
            kpi_counts[key][k] += cats[k]
    print(f"ｸﾛｰｻﾞｰ集計: {len(kpi_counts)}キー")

    # ③ イベントシート読み込み（A1:AA200）
    event_rows = execute_with_retry(sheets_svc.spreadsheets().values().get(
        spreadsheetId=SPREADSHEET_ID,
        range=f"'{event_sheet_name}'!A1:AA200"
    )).get("values", [])

    # ④ イベント日キーを構築（kpi_store 形式: "楽天経堂" etc.）
    event_kpi_keys: set = set()
    for i, row in enumerate(event_rows):
        if i < 3:
            continue
        store    = row[3].strip() if len(row) > 3 else ""
        date_str = row[4].strip() if len(row) > 4 else ""
        if not store or not date_str:
            continue
        try:
            kpi_store = STORE_TO_KPI.get(store)
            if kpi_store:
                event_kpi_keys.add((datetime.strptime(date_str, "%Y/%m/%d").date(), kpi_store))
        except ValueError:
            pass

    # ⑤ 戻り集計: イベント日以外のｸﾛｰｻﾞｰ実績を店舗ごとに合算（FILTER_YM適用）
    kpi_to_store = {v: k for k, v in STORE_TO_KPI.items()}
    kpi_return: dict = {}  # sheet_store → {mnp, hikari, turbo, card}
    for (d, kpi_store), cats in kpi_counts.items():
        if (d, kpi_store) in event_kpi_keys:
            continue
        if FILTER_YM and (d.year, d.month) != FILTER_YM:
            continue
        sheet_store = kpi_to_store.get(kpi_store)
        if not sheet_store:
            continue
        if sheet_store not in kpi_return:
            kpi_return[sheet_store] = {"mnp": 0, "hikari": 0, "turbo": 0, "card": 0}
        for k in cats:
            kpi_return[sheet_store][k] += cats[k]
    print(f"戻り集計: {len(kpi_return)}店舗")

    # ★ 書き込み前にクリア（再実行時の古いデータ・重複行を消す）
    execute_with_retry(sheets_svc.spreadsheets().values().batchClear(
        spreadsheetId=SPREADSHEET_ID,
        body={"ranges": [
            f"'{event_sheet_name}'!X4:AA200",          # 単日実績（イベント行）
            f"'{event_sheet_name}'!AD1:AH{4 + len(STORE_ORDER) + 2}",  # 戻りセクション全体
        ]}
    ))

    updates = []
    event_row_count = 0

    # ⑥ X〜AA列: イベント行ごとのカウントを書き込む
    for i, row in enumerate(event_rows):
        if i < 3:
            continue
        store    = row[3].strip() if len(row) > 3 else ""
        date_str = row[4].strip() if len(row) > 4 else ""
        if not store or not date_str:
            continue
        try:
            event_date = datetime.strptime(date_str, "%Y/%m/%d").date()
        except ValueError:
            continue
        kpi_store = STORE_TO_KPI.get(store)
        if not kpi_store:
            continue
        counts = kpi_counts.get((event_date, kpi_store), {"mnp": 0, "hikari": 0, "turbo": 0, "card": 0})
        row_num = i + 1
        updates.append({
            "range":  f"'{event_sheet_name}'!X{row_num}:AA{row_num}",
            "values": [[counts["mnp"], counts["hikari"], counts["turbo"], counts["card"]]],
        })
        event_row_count += 1
        print(f"  行{row_num}: {store} {event_date} → MNP={counts['mnp']}, 光={counts['hikari']}, Turbo={counts['turbo']}, カード={counts['card']}")

    # ⑦ AD〜AH列: 戻りセクション
    #   行1〜2: セクションタイトル（左側の単日実績タイトル行に合わせる）
    #   行3:    列ヘッダー（左側の黄色ヘッダー行に合わせる）
    #   行4〜:  店舗別合計（左側のデータ行に合わせる）
    ym_label = f"{FILTER_YM[0]}年{FILTER_YM[1]}月" if FILTER_YM else "当月"
    updates.append({
        "range":  f"'{event_sheet_name}'!AD1:AH1",
        "values": [[f"戻り（{ym_label} イベント日以外のクローザー獲得）", "", "", "", ""]],
    })
    updates.append({
        "range":  f"'{event_sheet_name}'!AD3:AH3",
        "values": [["店舗", "新規MNP", "ひかり", "Turbo", "クレカ"]],
    })
    for idx, store in enumerate(STORE_ORDER):
        row_num = 4 + idx
        c = kpi_return.get(store, {"mnp": 0, "hikari": 0, "turbo": 0, "card": 0})
        updates.append({
            "range":  f"'{event_sheet_name}'!AD{row_num}:AH{row_num}",
            "values": [[store, c["mnp"], c["hikari"], c["turbo"], c["card"]]],
        })
        print(f"  [戻り] {store} → MNP={c['mnp']}, 光={c['hikari']}, Turbo={c['turbo']}, カード={c['card']}")

    if not updates:
        print("更新なし")
        return

    execute_with_retry(sheets_svc.spreadsheets().values().batchUpdate(
        spreadsheetId=SPREADSHEET_ID,
        body={"valueInputOption": "RAW", "data": updates}
    ))
    print(f"=== KPI列入力完了: 単日実績 {event_row_count}行 / 戻り {len(STORE_ORDER)}店舗 ===")

# ===== メイン =====
def main():
    print(f"=== 同期開始 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} ===")

    # 必須環境変数チェック
    for _var in ("SPREADSHEET_ID", "KPI_SPREADSHEET_ID"):
        if not os.environ.get(_var):
            raise SystemExit(f"[ERROR] 環境変数 {_var} が設定されていません")

    gmail_creds  = build_creds("GMAIL")
    sheets_creds = build_creds("SHEETS")
    gmail_svc    = build("gmail",  "v1", credentials=gmail_creds)
    sheets_svc   = build("sheets", "v4", credentials=sheets_creds)

    processed = set() if RESET_PROCESSED else load_processed()
    if RESET_PROCESSED:
        print("★ RESET_PROCESSED=true: 処理済みキャッシュをリセット")
    if FILTER_YM:
        print(f"★ FILTER_YYYYMM: {FILTER_YM[0]}年{FILTER_YM[1]}月分のみ対象")
    sheet_name = get_current_sheet_name()
    print(f"対象シート: {sheet_name}")

    # シートが存在しなければテンプレートから自動作成（月初の create_monthly_sheet.yml が
    # 失敗した場合の自己修復フォールバック）
    sheet_id = ensure_current_sheet(sheets_svc, sheet_name)
    # 重複行クリーンアップ（過去の実行で残ったデータを除去）
    deduplicate_sheet(sheets_svc, sheet_name, sheet_id)
    row_map, date_rows = load_sheet_rows(sheets_svc, sheet_name)
    # 最初にデータが存在する最小行番号 = 書式（ドロップダウン）が整った基準行
    first_formatted_row = min(row_map.values()) if row_map else 4
    print(f"既存行数: {len(row_map)}  書式基準行: {first_formatted_row}")

    updated = 0

    for kind, query in [("起案", f"from:@{SENDER_DOMAIN} 起案"),
                        ("振り返り", f"from:@{SENDER_DOMAIN} 振り返り")]:

        # 振り返りは起案による行挿入後の最新状態が必要なため再読み込み
        if kind == "振り返り":
            row_map, date_rows = load_sheet_rows(sheets_svc, sheet_name)

        # ページネーションで全件取得（maxResults=50 の打ち切りを防ぐ）
        messages = []
        page_token = None
        while True:
            kwargs = dict(userId="me", q=query, maxResults=500)
            if page_token:
                kwargs["pageToken"] = page_token
            result = gmail_svc.users().messages().list(**kwargs).execute()
            messages.extend(result.get("messages", []))
            page_token = result.get("nextPageToken")
            if not page_token:
                break
        print(f"\nGmail '{query}': {len(messages)}件")

        for m in messages:
            if m["id"] in processed:
                continue

            msg = gmail_svc.users().messages().get(
                userId="me", id=m["id"], format="full"
            ).execute()
            headers = {h["name"]: h["value"] for h in msg["payload"]["headers"]}
            subject = headers.get("Subject", "")
            print(f"\n  [{kind}] {subject[:70]}")

            # 添付ファイルを読む
            att_text = get_attachment_text(gmail_svc, m["id"], msg["payload"])

            if att_text.strip():
                ok, store, dates, diff = verify(subject, att_text)
            else:
                # 添付なし → 件名のみ
                print("    [添付なし] 件名のみで処理")
                store = normalize_store(subject)
                dates = extract_dates(subject)
                ok, diff = True, ""

            if not ok:
                print(f"    [不一致] {diff}")
                if NOTIFY_EMAIL:
                    send_discrepancy_email(gmail_svc, msg, subject, diff)
                else:
                    print("    → 通知メール送信: スキップ（NOTIFY_EMAIL=false）")
                # 不一致でもシートへの入力は続行（内容が正しく再送不要なケースに対応）
                # store・dates は verify() が添付優先で返した値を使う

            if not store or not dates:
                print("    [WARN] 店舗または日付を特定できません")
                processed.add(m["id"])
                continue

            if kind == "起案":
                # 「再送」で始まる件名は古い行を削除してから挿入（店舗・日付が修正された場合も対応）
                is_resend = bool(re.match(r'^(再送\s*)+', subject.strip()))
                if is_resend:
                    target_dates = [d for d in dates
                                    if not FILTER_YM or (d.year, d.month) == FILTER_YM]
                    if delete_rows_by_dates(sheets_svc, sheet_name, sheet_id,
                                            target_dates, row_map):
                        # 行番号が変わっているので再読み込み
                        row_map, date_rows = load_sheet_rows(sheets_svc, sheet_name)

                for d in dates:
                    # 対象年月フィルタ
                    if FILTER_YM and (d.year, d.month) != FILTER_YM:
                        continue
                    date_rows, row_map = insert_event_row(
                        sheets_svc, sheet_name, store, d, date_rows, row_map,
                        sheet_id, first_formatted_row
                    )
                    updated += 1
            else:
                matched = False
                for d in dates:
                    # 対象年月フィルタ
                    if FILTER_YM and (d.year, d.month) != FILTER_YM:
                        continue
                    key = (store, d)
                    if key in row_map:
                        u = update_cell(sheets_svc, sheet_name, row_map[key], "振り返り")
                        print(f"    → {store} {d} 行{row_map[key]} 振り返り {'✓' if u else '(既存)'}")
                        if u:
                            updated += 1
                        matched = True
                if not matched:
                    print(f"    [行なし] {store} {dates} → 次回再試行")
                    continue

            processed.add(m["id"])

    save_processed(processed)
    print(f"\n=== 完了: {updated}件更新 ===")

    # KPI列（X〜AA）を更新
    fill_kpi_columns(sheets_svc, sheet_name)

if __name__ == "__main__":
    main()
